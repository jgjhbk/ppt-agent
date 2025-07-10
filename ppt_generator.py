import re
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from document_parser import DocumentModel  # 导入DocumentModel类

class PPTGenerator:
    """PPT生成器，负责将文档内容转换为PPT（每张幻灯片内容不超过9行）"""

    def __init__(self, template_path: Optional[str] = None):
        self.template_path = template_path
        self.prs = Presentation(template_path)
        self.slide_height = self.prs.slide_height
        self.body_shape_width = self.prs.slide_width - Inches(2.0)  # 文本框宽度
        self.max_lines_per_slide = 9  # 每页最大行数
        self.regular_font_size = Pt(18)
        self.code_font_size = Pt(14)
        self.line_spacing = 1.5

    def generate(self, doc_model: DocumentModel, output_path: str) -> None:
        """生成PPT文件"""
        # 创建标题页
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)

        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = doc_model.get_title() or "演示文稿"
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "宋体"
                run.font.size = Pt(40)  
                run.font.bold = True
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.line_spacing = 1.5

        subtitle.text = f"作者: {doc_model.get_author() or '未知'}\n日期: {doc_model.metadata.get('date', '')}"
        for paragraph in subtitle.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "宋体"
                run.font.bold = True
                run.font.size = Pt(20)  
            paragraph.line_spacing = 1.5

        # 创建内容页
        self._process_sections(doc_model.sections, doc_model.images)

        self.prs.save(output_path)

    def _process_sections(self, sections: List[Dict], images: Dict[str, Any], parent_level: int = 0) -> None:
        """递归处理章节，生成对应的幻灯片"""
        for section in sections:
            if section['level'] == 1:  
                self._add_title_only_slide(section['title'])
                
                if section['content'].strip():
                    self._add_content_slide(section['content'], images)

                if section.get('subsections'):
                    self._process_sections(section['subsections'], images, section['level'])
            
            elif section['level'] == 2:  # 二级标题
                self._add_bullet_slide(section['title'], section['content'], images)
                if section.get('subsections'):
                    self._process_sections(section['subsections'], images, section['level'])

    def _add_title_only_slide(self, title: str) -> None:
        """添加仅包含标题的幻灯片"""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)

        left = Inches(1.0)
        top = Inches(3.0)
        width = Inches(8.0)
        height = Inches(2.0)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.add_paragraph()
        p.text = title
        p.alignment = PP_ALIGN.CENTER

        for run in p.runs:
            run.font.name = "宋体"
            run.font.size = Pt(32)
            run.font.bold = True

    def _add_content_slide(self, content: str, images: Dict[str, Any]) -> None:
        """添加仅包含内容的幻灯片（用于一级标题的正文）"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)

        # 清除标题占位符的默认文本
        title_shape = slide.shapes.title
        title_shape.text = ""
        
        body_shape = slide.placeholders[1]
        self._populate_content(body_shape, content, images)

    def _add_bullet_slide(self, title: str, content: str, images: Dict[str, Any]) -> None:
        """添加带要点的幻灯片"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "宋体"
                run.font.size = Pt(28)
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.line_spacing = 1.5

        body_shape = slide.placeholders[1]
        self._populate_content(body_shape, content, images)

    def _populate_content(self, body_shape, content: str, images: Dict[str, Any]) -> None:
        """处理内容，支持列表、图片、代码块和&&分页标记"""
        tf = body_shape.text_frame
        tf.text = ""

        if not content:
            p = tf.add_paragraph()
            p.text = " "
            self._set_font_and_spacing(p)
            return

        # 检测分页标记&&，并分割内容
        page_break_pattern = re.compile(r'&&')
        content_parts = page_break_pattern.split(content)
        
        # 处理第一部分内容（当前幻灯片）
        first_part = content_parts[0].strip()
        if first_part:
            self._process_single_part(tf, first_part, images)
        
        # 处理剩余部分（创建新幻灯片）
        for part in content_parts[1:]:
            part = part.strip()
            if part:
                # 创建新幻灯片
                new_slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
                new_slide.shapes.title.text = ""  # 清空标题
                # 处理新幻灯片内容
                new_body = new_slide.placeholders[1]
                self._process_single_part(new_body.text_frame, part, images)

    def _process_single_part(self, text_frame, content: str, images: Dict[str, Any]) -> int:
        """处理单部分内容（不含分页标记的文本块）
        返回已使用的行数
        """
        image_pattern = re.compile(r'!\[(.*?)\]\((.*?)\)')
        code_block_pattern = re.compile(r'```([\s\S]*?)```')
        
        code_blocks = list(code_block_pattern.finditer(content))
        last_end = 0
        used_lines = 0  # 记录已使用的行数
        
        for code_match in code_blocks:
            pre_code_text = content[last_end:code_match.start()].strip()
            if pre_code_text:
                # 处理代码块前的文本
                used_lines = self._process_regular_content(text_frame, pre_code_text, images, used_lines)

            # 处理代码块
            code_content = code_match.group(1).strip()
            
            # 估计代码块的行数
            code_lines = self._estimate_text_lines(code_content, self.code_font_size)
            
            # 如果当前幻灯片放不下代码块，创建新幻灯片
            if used_lines + code_lines > self.max_lines_per_slide:
                text_frame = self._create_new_text_slide().placeholders[1].text_frame
                used_lines = 0
            
            # 添加代码块内容
            p = text_frame.add_paragraph()
            p.text = code_content
            p.level = 0
            p.alignment = PP_ALIGN.LEFT
            
            # 设置代码字体
            for run in p.runs:
                run.font.name = "Consolas"
                run.font.size = Pt(12)
            p.line_spacing = self.line_spacing
            
            used_lines += code_lines
            last_end = code_match.end()

        # 处理所有代码块之后的文本
        remaining_text = content[last_end:].strip()
        if remaining_text:
            used_lines = self._process_regular_content(text_frame, remaining_text, images, used_lines)
            
        return used_lines

    def _process_regular_content(self, text_frame, content: str, images: Dict[str, Any], used_lines: int) -> int:
        """处理常规内容，添加自动分页逻辑"""
        paragraphs = re.split(r'(\n\n)', content)
        paragraphs = [p for p in paragraphs if p.strip()]
        
        current_text_frame = text_frame
        current_used_lines = used_lines
        
        for paragraph in paragraphs:
            if not paragraph.strip():
                continue
                
            image_match = re.search(r'!\[(.*?)\]\((.*?)\)', paragraph)
            if image_match:
                alt_text = image_match.group(1)
                url = image_match.group(2)
                
                image_id = None
                for img_id, img_info in images.items():
                    if img_info['url'] == url:
                        image_id = img_id
                        break
                
                if image_id and images[image_id]['local_path']:
                    # 创建新幻灯片放图片
                    new_slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
                    new_slide.shapes.title.text = ""  # 清空标题
                    self._add_image_slide(new_slide, alt_text, images[image_id]['local_path'])
                    # 创建下一张文本幻灯片
                    current_text_frame = self._create_new_text_slide().placeholders[1].text_frame
                    current_used_lines = 0  # 重置行数计数
                else:
                    p = current_text_frame.add_paragraph()
                    p.text = f"[图片: {alt_text} - 无法加载]"
                    p.level = 0
                    p.alignment = PP_ALIGN.LEFT
                    p.font.size = self.regular_font_size
                    self._set_font_and_spacing(p)
                    
                    # 估计这一段的行数
                    lines_added = self._estimate_text_lines(p.text, self.regular_font_size)
                    current_used_lines += lines_added
                    
                    # 检查是否需要分页
                    if current_used_lines >= self.max_lines_per_slide:
                        current_text_frame = self._create_new_text_slide().placeholders[1].text_frame
                        current_used_lines = 0
            else:
                lines = paragraph.strip().split('\n')
                if not lines:
                    continue

                if lines[0].startswith('- '):
                    # 处理列表内容
                    line_count = len(lines)
                    
                    if current_used_lines + line_count > self.max_lines_per_slide:
                        # 需要分页
                        remaining_lines = lines
                        while remaining_lines:
                            lines_to_add = remaining_lines[:self.max_lines_per_slide - current_used_lines]
                            remaining_lines = remaining_lines[self.max_lines_per_slide - current_used_lines:]
                            
                            self._add_list(current_text_frame, lines_to_add)
                            current_used_lines += len(lines_to_add)
                            
                            if remaining_lines:
                                current_text_frame = self._create_new_text_slide().placeholders[1].text_frame
                                current_used_lines = 0
                    else:
                        self._add_list(current_text_frame, lines)
                        current_used_lines += line_count
                else:
                    # 处理普通段落
                    # 估计这一段的行数
                    lines_added = self._estimate_text_lines(paragraph, self.regular_font_size)
                    
                    if current_used_lines + lines_added > self.max_lines_per_slide:
                        # 需要分页
                        current_text_frame = self._create_new_text_slide().placeholders[1].text_frame
                        current_used_lines = 0
                    
                    p = current_text_frame.add_paragraph()
                    p.text = paragraph
                    p.level = 0
                    p.alignment = PP_ALIGN.LEFT
                    p.font.size = self.regular_font_size
                    self._set_font_and_spacing(p)
                    
                    current_used_lines += lines_added

        return current_used_lines

    def _create_new_text_slide(self):
        """创建一个新的文本幻灯片"""
        new_slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        new_slide.shapes.title.text = ""  # 清空标题
        return new_slide

    def _add_list(self, text_frame, lines: List[str]) -> None:
        """添加列表内容"""
        for line in lines:
            line = line.rstrip()
            if not line:
                continue

            level = 0
            while line.startswith('  '):
                level += 1
                line = line[2:].lstrip()

            if line.startswith('- '):
                p = text_frame.add_paragraph()
                p.text = line[2:].strip()
                p.level = min(level, 4)
                p.alignment = PP_ALIGN.LEFT
                p.font.size = Pt(18 - level * 2)
                self._set_font_and_spacing(p)
            else:
                p = text_frame.add_paragraph()
                p.text = line
                p.level = 0
                p.alignment = PP_ALIGN.LEFT
                p.font.size = self.regular_font_size
                self._set_font_and_spacing(p)

    def _set_font_and_spacing(self, paragraph) -> None:
        """设置字体和行距"""
        for run in paragraph.runs:
            run.font.name = "宋体"
        paragraph.line_spacing = self.line_spacing

    def _estimate_text_lines(self, text, font_size):
        """估计文本在给定字体大小下的行数"""
        # 简化的估计算法
        avg_char_width = font_size * 0.5  # 假设平均字符宽度
        chars_per_line = self.body_shape_width / avg_char_width
        
        # 空行也需要计算
        lines = text.split('\n')
        estimated_lines = 0
        
        for line in lines:
            # 每行至少占一行
            estimated_lines += 1
            
            # 非空行需要计算可能的折行
            if line.strip() != '':
                estimated_lines += max(0, len(line) / chars_per_line - 1)
        
        return estimated_lines

    def _add_code_block(self, text_frame, code_content: str, used_lines: int) -> int:
        """添加代码块并返回新增的行数"""
        # 估计代码块的行数
        estimated_lines = self._estimate_text_lines(code_content, self.code_font_size)
        
        # 检查是否需要分页
        if used_lines + estimated_lines > self.max_lines_per_slide:
            # 需要分页
            text_frame = self._create_new_text_slide().placeholders[1].text_frame
            used_lines = 0
        
        # 添加代码块
        p = text_frame.add_paragraph()
        p.text = code_content
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        
        # 设置代码字体
        for run in p.runs:
            run.font.name = "Consolas"
            run.font.size = self.code_font_size
        p.line_spacing = self.line_spacing
        
        return used_lines + estimated_lines

    def _add_image_slide(self, slide, alt_text: str, image_path: str) -> None:
        """添加图片幻灯片"""
        title = slide.shapes.title
        title.text = alt_text
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "宋体"
                run.font.size = Pt(28)
            paragraph.line_spacing = 1.5

        try:
            # 尝试获取图片尺寸
            from PIL import Image
            img = Image.open(image_path)
            width, height = img.size
            
            # 计算图片在幻灯片中的最佳尺寸
            slide_width = self.prs.slide_width
            slide_height = self.prs.slide_height - Inches(1.5)  # 留出标题空间
            
            # 保持图片比例
            img_ratio = width / height
            slide_ratio = slide_width / slide_height
            
            if img_ratio > slide_ratio:
                # 图片更宽，以宽度为基准
                img_width = slide_width - Inches(1.0)
                img_height = img_width / img_ratio
            else:
                # 图片更高，以高度为基准
                img_height = slide_height - Inches(1.0)
                img_width = img_height * img_ratio
            
            # 居中放置图片
            left = (slide_width - img_width) / 2
            top = Inches(1.5)  # 标题下方的位置
            
            slide.shapes.add_picture(image_path, left, top, width=img_width, height=img_height)
        except Exception as e:
            # 如果出现错误，添加文本说明
            print(f"无法添加图片 {image_path}: {str(e)}")
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            p = tf.add_paragraph()
            p.text = f"无法加载图片: {alt_text}"
            p.font.name = "宋体"
            p.font.size = self.regular_font_size    
